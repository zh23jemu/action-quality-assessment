import json
import math
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path.cwd()
DOCS = ROOT / 'docs'
ASSETS = DOCS / 'ppt_assets'
DOCS.mkdir(exist_ok=True)
ASSETS.mkdir(exist_ok=True)

# -----------------------------
# 读取实验结果：所有数字来自服务器同步回来的正式 JSON。
# -----------------------------
metrics_dir = ROOT / 'outputs' / 'metrics'
preds_dir = ROOT / 'outputs' / 'predictions'
motion_train = json.loads((metrics_dir / 'motion_train.json').read_text(encoding='utf-8'))
motion_eval = json.loads((metrics_dir / 'motion_eval_real.json').read_text(encoding='utf-8'))
pose_train = json.loads((metrics_dir / 'pose_train.json').read_text(encoding='utf-8'))
pose_eval = json.loads((metrics_dir / 'pose_eval_real.json').read_text(encoding='utf-8'))
motion_pred = json.loads((preds_dir / 'motion_single_real.json').read_text(encoding='utf-8'))
pose_pred = json.loads((preds_dir / 'pose_single_real.json').read_text(encoding='utf-8'))

# -----------------------------
# 视觉常量：克制的技术汇报风格，避免照搬参考图里的内容。
# -----------------------------
NAVY = RGBColor(18, 38, 74)
BLUE = RGBColor(38, 93, 170)
TEAL = RGBColor(0, 142, 150)
CYAN = RGBColor(40, 178, 190)
GREEN = RGBColor(32, 151, 96)
ORANGE = RGBColor(238, 132, 54)
RED = RGBColor(207, 75, 75)
INK = RGBColor(28, 36, 48)
MUTED = RGBColor(91, 103, 120)
LIGHT = RGBColor(246, 249, 252)
LINE = RGBColor(220, 228, 238)
WHITE = RGBColor(255, 255, 255)

TITLE_FONT = 'Microsoft YaHei'
BODY_FONT = 'Microsoft YaHei'

plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'Arial Unicode MS', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# -----------------------------
# 图表资产
# -----------------------------
def save_motion_curve():
    hist = motion_train['history']
    epochs = [h['epoch'] for h in hist]
    src = [h['SRC'] for h in hist]
    rl2 = [h['R-L2'] for h in hist]
    loss = [h['train_loss'] for h in hist]
    fig, axes = plt.subplots(1, 3, figsize=(13.5, 3.5), dpi=180)
    axes[0].plot(epochs, loss, color='#265DAA', linewidth=2.5)
    axes[0].set_title('训练 Loss 下降')
    axes[0].set_xlabel('Epoch')
    axes[0].set_ylabel('MSE')
    axes[1].plot(epochs, src, color='#008E96', linewidth=2.5)
    axes[1].scatter([30], [motion_eval['SRC']], color='#EE8436', zorder=3)
    axes[1].set_title('SRC 持续提升')
    axes[1].set_xlabel('Epoch')
    axes[1].set_ylim(0.2, 0.82)
    axes[2].plot(epochs, rl2, color='#CF4B4B', linewidth=2.5)
    axes[2].scatter([30], [motion_eval['R-L2']], color='#EE8436', zorder=3)
    axes[2].set_title('R-L2 降低')
    axes[2].set_xlabel('Epoch')
    for ax in axes:
        ax.grid(True, color='#DCE4EE', linewidth=0.8)
        ax.spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    out = ASSETS / 'motion_curves.png'
    fig.savefig(out, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return out


def save_pose_curve():
    hist = pose_train['history']
    epochs = [h['epoch'] for h in hist]
    mean_f1 = [h['mean_f1'] for h in hist]
    loss = [h['train_loss'] for h in hist]
    fig, axes = plt.subplots(1, 2, figsize=(10, 3.6), dpi=180)
    axes[0].plot(epochs, loss, color='#265DAA', linewidth=2.5)
    axes[0].set_title('训练 Loss 下降')
    axes[0].set_xlabel('Epoch')
    axes[0].grid(True, color='#DCE4EE')
    axes[0].spines[['top', 'right']].set_visible(False)
    axes[1].plot(epochs, mean_f1, color='#008E96', linewidth=2.5)
    axes[1].scatter([25], [pose_eval['mean_f1']], color='#EE8436', zorder=3)
    axes[1].set_title('Mean F1 最佳出现在第 25 轮')
    axes[1].set_xlabel('Epoch')
    axes[1].grid(True, color='#DCE4EE')
    axes[1].spines[['top', 'right']].set_visible(False)
    fig.tight_layout()
    out = ASSETS / 'pose_curves.png'
    fig.savefig(out, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return out


def save_pose_bar():
    f1 = pose_eval['F1']
    names = ['position', 'armstand', 'rotation_type', 'somersaults', 'twists']
    vals = [f1[n] for n in names]
    fig, ax = plt.subplots(figsize=(8.5, 3.4), dpi=180)
    bars = ax.bar(names, vals, color=['#265DAA', '#008E96', '#25B2BE', '#EE8436', '#CF4B4B'])
    ax.axhline(pose_eval['mean_f1'], color='#12264A', linestyle='--', linewidth=1.5, label=f"mean={pose_eval['mean_f1']:.3f}")
    ax.set_ylim(0, 0.65)
    ax.set_ylabel('Macro F1')
    ax.set_title('Pose Contrastive 各属性 F1')
    ax.grid(axis='y', color='#DCE4EE')
    ax.spines[['top', 'right']].set_visible(False)
    ax.legend(frameon=False)
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x()+bar.get_width()/2, val+0.015, f'{val:.3f}', ha='center', fontsize=9)
    fig.tight_layout()
    out = ASSETS / 'pose_f1_bar.png'
    fig.savefig(out, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return out

motion_curve = save_motion_curve()
pose_curve = save_pose_curve()
pose_bar = save_pose_bar()

# -----------------------------
# PPT 工具函数
# -----------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def rgb_tuple(c):
    return (c[0], c[1], c[2])


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=LINE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_bg(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, color)
    bg.line.fill.background()
    return bg


def add_accent(slide, color=TEAL):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height)
    set_fill(bar, color)
    bar.line.fill.background()


def add_textbox(slide, text_value, x, y, w, h, font_size=22, color=INK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=BODY_FONT, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text_value
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, font_size=20, color=INK, bullet_color=TEAL, gap=0.42):
    # 使用独立文本框做项目符号，避免 PowerPoint 自动缩进在中文下变形。
    for i, item in enumerate(items):
        yy = y + i * gap
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy+0.09), Inches(0.1), Inches(0.1))
        set_fill(dot, bullet_color)
        dot.line.fill.background()
        add_textbox(slide, item, x+0.22, yy, w-0.22, gap, font_size=font_size, color=color)


def add_title(slide, section, title, subtitle=None):
    add_accent(slide)
    add_textbox(slide, section, 0.65, 0.35, 5.4, 0.32, font_size=15, color=TEAL, bold=True)
    add_textbox(slide, title, 0.65, 0.82, 11.8, 0.62, font_size=34, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, 0.68, 1.35, 11.5, 0.42, font_size=16, color=MUTED)


def add_card(slide, x, y, w, h, title=None, body=None, fill_color=WHITE, title_color=NAVY, border=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill_color)
    set_line(shape, border, 1)
    if title:
        add_textbox(slide, title, x+0.25, y+0.2, w-0.5, 0.35, font_size=18, color=title_color, bold=True)
    if body:
        add_textbox(slide, body, x+0.25, y+0.68, w-0.5, h-0.85, font_size=15, color=INK)
    return shape


def add_metric(slide, x, y, w, h, label, value, note='', color=BLUE):
    add_card(slide, x, y, w, h, fill_color=WHITE)
    add_textbox(slide, label, x+0.25, y+0.22, w-0.5, 0.35, font_size=15, color=MUTED, bold=True)
    add_textbox(slide, value, x+0.25, y+0.65, w-0.5, 0.58, font_size=32, color=color, bold=True)
    if note:
        add_textbox(slide, note, x+0.25, y+1.24, w-0.5, 0.35, font_size=12, color=MUTED)


def add_footer(slide, idx):
    add_textbox(slide, f'MTL-AQA × Fitness-AQA 适配复现 | {idx:02d}', 10.5, 7.05, 2.4, 0.24, font_size=9, color=RGBColor(130,140,155), align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, columns, rows, font_size=12):
    table = slide.shapes.add_table(len(rows)+1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h)).table
    col_w = w / len(columns)
    for i, col in enumerate(columns):
        table.columns[i].width = Inches(col_w)
        cell = table.cell(0, i)
        set_fill(cell, NAVY)
        cell.text = col
        for p in cell.text_frame.paragraphs:
            p.font.name = BODY_FONT
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            set_fill(cell, WHITE if r % 2 else RGBColor(242, 247, 250))
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.name = BODY_FONT
                p.font.size = Pt(font_size)
                p.font.color.rgb = INK
    return table

# -----------------------------
# Slide 1: Cover
# -----------------------------
slide = prs.slides.add_slide(blank)
add_bg(slide, NAVY)
# 大块色带和视觉节奏
shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.8), Inches(0), Inches(4.6), Inches(7.5))
set_fill(shape, RGBColor(20, 73, 120)); shape.line.fill.background()
shape.rotation = 0
add_textbox(slide, 'MTL-AQA 上的动作质量评估复现', 0.8, 0.95, 8.7, 0.72, font_size=38, color=WHITE, bold=True)
add_textbox(slide, 'Motion Disentangling × Pose Contrastive', 0.82, 1.75, 8.2, 0.42, font_size=20, color=RGBColor(188, 225, 230), bold=True)
add_textbox(slide, '从数据下载、抽帧、特征生成，到训练、测试指标与单样本推理的端到端工程闭环', 0.84, 2.45, 7.6, 0.85, font_size=20, color=RGBColor(230, 241, 246))
add_metric(slide, 0.85, 4.35, 2.55, 1.35, '数据规模', '1412', '1059 train / 353 test', color=CYAN)
add_metric(slide, 3.72, 4.35, 2.55, 1.35, 'Motion SRC', '0.7661', 'R-L2 = 0.1683', color=GREEN)
add_metric(slide, 6.58, 4.35, 2.55, 1.35, 'Pose mean F1', '0.4052', '5 类动作属性', color=ORANGE)
add_textbox(slide, '第 18 组 | 课程实践汇报', 0.85, 6.65, 4.5, 0.35, font_size=15, color=RGBColor(188, 225, 230))

# Slide 2: Goals
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '01 任务目标', '客户需求拆解：两个模型、一个数据集、三类结果')
add_card(slide, 0.8, 1.95, 3.75, 3.65, '指定数据集', 'MTL-AQA 跳水动作质量评估数据集\n\n15 个完整比赛视频中标注出 1412 个有效跳水片段。')
add_card(slide, 4.8, 1.95, 3.75, 3.65, '两个模型', 'Motion Disentangling：质量分数回归\n\nPose Contrastive：跳水动作属性分类')
add_card(slide, 8.8, 1.95, 3.75, 3.65, '要求指标', 'Motion 输出 SRC 与 R-L2\n\nPose 输出 F1 Score\n\n另需单样本推理与问题记录')
add_footer(slide, 2)

# Slide 3: Environment
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '02 实验环境', '在 Slurm GPU 节点完成真实训练和复验')
add_metric(slide, 0.8, 1.9, 2.5, 1.15, 'GPU', 'NVIDIA L40S', 'Slurm 计算节点', color=BLUE)
add_metric(slide, 3.55, 1.9, 2.5, 1.15, 'CUDA', '12.6', 'torch CUDA runtime', color=TEAL)
add_metric(slide, 6.3, 1.9, 2.5, 1.15, 'PyTorch', '2.12.0', '+cu126', color=GREEN)
add_metric(slide, 9.05, 1.9, 2.5, 1.15, '调度', 'Slurm', 'aws / gpo-ifv7xx', color=ORANGE)
add_card(slide, 0.8, 3.55, 5.8, 2.5, '环境验证', '服务器上 `torch.cuda.is_available()` 为 True，GPU 识别为 NVIDIA L40S。\n\n训练和 eval 均记录了相同的 CUDA / GPU 环境摘要。')
add_card(slide, 7.0, 3.55, 5.4, 2.5, '工程约定', '所有 Python 入口使用项目 `.venv`。\n\n长时间训练通过 Slurm 脚本提交；日志、指标、checkpoint、预测结果统一保存到 `outputs/`。')
add_footer(slide, 3)

# Slide 4: Dataset
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '03 数据预处理', '不是训练 15 条长视频，而是训练官方标注的动作片段')
add_metric(slide, 0.75, 1.75, 2.45, 1.2, '完整视频', '15 个', 'YouTube / RTIS 来源', color=BLUE)
add_metric(slide, 3.45, 1.75, 2.45, 1.2, '有效片段', '1412 个', '官方 start/end frame', color=TEAL)
add_metric(slide, 6.15, 1.75, 2.45, 1.2, '训练集', '1059', 'split 0', color=GREEN)
add_metric(slide, 8.85, 1.75, 2.45, 1.2, '测试集', '353', 'split 0', color=ORANGE)
# 流程图
steps = [('下载视频', '15 个长视频'), ('全量抽帧', '建立帧库'), ('读取标注', 'start_frame/end_frame'), ('片段采样', '只取有效动作'), ('生成特征', '128 维帧统计特征')]
for i, (t, b) in enumerate(steps):
    x = 0.75 + i*2.45
    add_card(slide, x, 4.05, 1.95, 1.25, t, b, fill_color=WHITE)
    if i < len(steps)-1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+1.85), Inches(4.43), Inches(0.45), Inches(0.28))
        set_fill(arr, TEAL); arr.line.fill.background()
add_textbox(slide, '关键回答：视频虽然 1-2 小时，但模型只看到官方标注出的跳水动作片段帧，不使用等待、回放、裁判和转场等无关画面。', 0.9, 6.1, 11.5, 0.55, font_size=17, color=NAVY, bold=True)
add_footer(slide, 4)

# Slide 5: Model overview
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '04 模型与适配实现', '官方仓库不完整，因此补齐可运行工程链路')
add_card(slide, 0.78, 1.75, 5.65, 3.6, 'Motion Disentangling', '目标：预测动作质量分数\n\n输入：片段级视频特征\n\n训练：回归损失\n\n指标：SRC 衡量排序一致性，R-L2 衡量相对误差')
add_card(slide, 6.9, 1.75, 5.65, 3.6, 'Pose Contrastive', '目标：学习动作属性表征\n\n输入：片段级视频特征\n\n训练：多属性分类 / 表征适配\n\n指标：position、armstand、rotation_type、somersaults、twists 的 F1')
add_textbox(slide, '说明：公开 GitHub 仓库无法原封不动跑出 MTL-AQA 上的完整训练测试；本项目根据论文思路补齐数据加载、特征、训练、评估和推理入口。', 1.0, 6.05, 11.2, 0.65, font_size=17, color=RED, bold=True)
add_footer(slide, 5)

# Slide 6: End-to-end project
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '05 工程闭环', '从数据到结果的统一入口已经补齐')
rows = [
    ['数据整理', 'scripts/prepare_data.py', '生成 manifest，校验 split / 标签 / 帧路径'],
    ['特征抽取', 'scripts/extract_features.py', '从有效动作片段采样帧，输出 feature_path'],
    ['Motion 训练/测试', 'train_motion.py / eval_motion.py', '输出 SRC、R-L2 与 checkpoint'],
    ['Pose 训练/测试', 'train_pose.py / eval_pose.py', '输出各属性 F1 与 mean F1'],
    ['单样本推理', 'scripts/infer_single.py', '输出分数或动作属性预测 JSON'],
    ['Slurm', 'slurm/*.slurm', 'GPU/CPU 作业、日志、指标统一保存'],
]
add_table(slide, 0.72, 1.72, 11.9, 4.8, ['模块', '入口', '作用'], rows, font_size=12)
add_footer(slide, 6)

# Slide 7: Key bugs
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '06 关键问题与修复', '这些问题适合放在 PPT 的“过程记录”页')
bugs = [
    ('数据下载', 'YouTube 网络与 Cookie 风控，服务器直连和代理均不稳定', '本地代理下载后通过 GitHub Release 同步到服务器'),
    ('视频解码', '03/06/18 为 AV1 编码，集群 ffmpeg 无法软解', '安装 imageio-ffmpeg，使用带 libdav1d 的新版 ffmpeg 补抽帧'),
    ('特征缺失', '初次特征抽取 feature_ready=1207，missing=205', '补抽 03/06/18 后 feature_ready=1412，missing=0'),
    ('路径解析', 'feature_path 被拼成 data/processed/data/processed/...', '修复 resolve_path，兼容项目根目录与 manifest 目录两种相对路径'),
    ('参数口径', 'eval 脚本原先只支持 --metrics，不支持 --output', '增加 --output 作为 --metrics 别名，命令更统一'),
]
add_table(slide, 0.62, 1.62, 12.1, 4.95, ['问题', '现象', '解决方案'], bugs, font_size=10.5)
add_textbox(slide, '价值：这些不是“失败”，而是工程复现过程中真实遇到并解决的兼容性和数据链路问题。', 0.85, 6.72, 11.3, 0.35, font_size=15, color=NAVY, bold=True)
add_footer(slide, 7)

# Slide 8: Motion results
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '07 Motion Disentangling 结果', '测试集 353 个片段，独立 eval 与训练期最佳指标一致')
add_metric(slide, 0.75, 1.65, 2.55, 1.22, 'SRC ↑', f"{motion_eval['SRC']:.4f}", 'Spearman 排序相关', color=GREEN)
add_metric(slide, 3.55, 1.65, 2.55, 1.22, 'R-L2 ↓', f"{motion_eval['R-L2']:.4f}", '相对 L2 距离', color=ORANGE)
add_metric(slide, 6.35, 1.65, 2.55, 1.22, 'Best epoch', '30', 'checkpoint 可复验', color=BLUE)
add_metric(slide, 9.15, 1.65, 2.55, 1.22, 'Test clips', '353', 'MTL-AQA split 0', color=TEAL)
slide.shapes.add_picture(str(motion_curve), Inches(0.75), Inches(3.2), width=Inches(11.9), height=Inches(3.1))
add_footer(slide, 8)

# Slide 9: Pose results
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '08 Pose Contrastive 结果', '输出跳水动作属性 F1，平均 F1 作为综合指标')
add_metric(slide, 0.75, 1.55, 2.55, 1.18, 'Mean F1 ↑', f"{pose_eval['mean_f1']:.4f}", '5 个属性平均', color=GREEN)
add_metric(slide, 3.55, 1.55, 2.55, 1.18, 'Best epoch', '25', 'checkpoint 可复验', color=BLUE)
add_metric(slide, 6.35, 1.55, 2.55, 1.18, 'Test clips', '353', 'MTL-AQA split 0', color=TEAL)
add_metric(slide, 9.15, 1.55, 2.55, 1.18, '属性数', '5 类', 'position / twists 等', color=ORANGE)
slide.shapes.add_picture(str(pose_curve), Inches(0.75), Inches(3.15), width=Inches(5.7), height=Inches(2.55))
slide.shapes.add_picture(str(pose_bar), Inches(6.65), Inches(3.05), width=Inches(5.6), height=Inches(2.65))
add_footer(slide, 9)

# Slide 10: Eval reproducibility
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '09 独立测试验证', '训练后 checkpoint 单独 eval，指标完全对齐')
rows = [
    ['Motion', '训练期最佳', f"SRC {motion_train['best_SRC']:.4f}", 'R-L2 0.1683', 'epoch 30'],
    ['Motion', '独立 eval', f"SRC {motion_eval['SRC']:.4f}", f"R-L2 {motion_eval['R-L2']:.4f}", '353 clips'],
    ['Pose', '训练期最佳', f"mean F1 {pose_train['best_mean_f1']:.4f}", '属性 F1 全量输出', 'epoch 25'],
    ['Pose', '独立 eval', f"mean F1 {pose_eval['mean_f1']:.4f}", '属性 F1 全量输出', '353 clips'],
]
add_table(slide, 0.75, 1.65, 11.85, 2.65, ['模型', '来源', '主指标', '辅助信息', '说明'], rows, font_size=12)
add_card(slide, 0.95, 4.75, 5.4, 1.35, '结论', 'checkpoint 可独立复验，不依赖训练过程中的临时状态。')
add_card(slide, 6.75, 4.75, 5.4, 1.35, '结果文件', 'motion_eval_real.json\npose_eval_real.json')
add_footer(slide, 10)

# Slide 11: Single inference
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '10 单样本推理', '用测试样本 01_0001 展示模型可部署入口')
add_card(slide, 0.8, 1.8, 5.6, 3.0, 'Motion 单样本输出', f"sample_id: {motion_pred['sample_id']}\n\n真实分数: {motion_pred['true_score']:.1f}\n预测分数: {motion_pred['pred_score']:.4f}\n\ncheckpoint: motion_best.pt")
attrs = pose_pred['predicted_attributes']
pose_body = '\n'.join([f'{k}: {v}' for k, v in attrs.items()])
add_card(slide, 6.9, 1.8, 5.6, 3.0, 'Pose 单样本输出', f"sample_id: {pose_pred['sample_id']}\n\n{pose_body}\n\ncheckpoint: pose_best.pt")
add_textbox(slide, '用途：这页可以配合终端截图，证明训练后的权重能够被独立加载并完成真实样本推理。', 1.0, 5.65, 11.1, 0.48, font_size=17, color=NAVY, bold=True)
add_footer(slide, 11)

# Slide 12: Gap analysis
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '11 与“完整官方复现”的差距', '公开仓库支撑可运行适配复现，但不支撑原封不动复现')
add_card(slide, 0.8, 1.75, 3.65, 3.55, '公开仓库缺口', 'MTL-AQA 仓库主要给数据和标注。\n\nFitness-AQA 仓库不包含可直接迁移到 MTL-AQA 的完整训练工程。')
add_card(slide, 4.85, 1.75, 3.65, 3.55, '特征链路差异', '论文级复现通常需要官方 C3D/I3D/pose 特征、预训练权重和完整配置。\n\n当前使用真实抽帧后的轻量帧统计特征。')
add_card(slide, 8.9, 1.75, 3.65, 3.55, '当前定位', '端到端可运行：数据、训练、测试、checkpoint、推理全部闭环。\n\n对外表述为“可运行适配复现”。')
add_textbox(slide, '建议口径：能跑通流程、能输出指标、能复验结果；若要求官方原始特征和源码级复现，需要作者提供更完整工程与实验配置。', 0.95, 6.05, 11.3, 0.55, font_size=16, color=RED, bold=True)
add_footer(slide, 12)

# Slide 13: Deliverables
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, '12 交付物状态', '程序侧已闭环，文档侧进入整理阶段')
rows = [
    ['代码工程', '已完成', 'scripts/、slurm/、README、AGENTS'],
    ['训练测试结果', '已完成', 'outputs/logs、outputs/metrics、outputs/checkpoints'],
    ['单样本推理', '已完成', 'outputs/predictions/*_single_real.json'],
    ['问题记录', '已形成素材', '下载、抽帧、路径、eval 参数等'],
    ['训练说明文档', '待整理', '按模板补环境、命令、截图'],
    ['课程总结报告', '待整理', '3000-5000 字，说明适配复现口径'],
]
add_table(slide, 0.8, 1.75, 11.75, 4.8, ['交付项', '状态', '说明'], rows, font_size=12)
add_footer(slide, 13)

# Slide 14: Summary
slide = prs.slides.add_slide(blank)
add_bg(slide, NAVY)
add_textbox(slide, '13 总结', 0.78, 0.65, 1.6, 0.35, font_size=18, color=CYAN, bold=True)
add_textbox(slide, '我们已经完成 MTL-AQA 上的端到端可运行适配复现', 0.78, 1.08, 11.2, 0.72, font_size=32, color=WHITE, bold=True)
summary = [
    '数据：15 个长视频中标注出的 1412 个跳水片段，特征缺失已清零。',
    '模型：Motion 与 Pose 两条链路均完成训练、独立测试、checkpoint 和单样本推理。',
    '指标：Motion SRC=0.7661、R-L2=0.1683；Pose mean F1=0.4052。',
    '工程：补齐了官方仓库缺失的数据加载、训练、评估、推理和 Slurm 脚本。',
    '边界：不是官方原始源码/原始特征的完整复现，需要在报告中说明。',
]
add_bullets(slide, summary, 0.95, 2.2, 11.5, 2.7, font_size=20, color=RGBColor(235,244,248), bullet_color=CYAN, gap=0.5)
add_card(slide, 0.95, 5.48, 11.1, 1.25, '下一步', '把训练日志截图、独立 eval 指标、单样本推理输出和关键问题修复整理进训练说明文档与课程报告。', fill_color=RGBColor(232, 246, 247), title_color=NAVY)

out = DOCS / 'Motion-Pose_汇报.pptx'
prs.save(out)
print(out)
